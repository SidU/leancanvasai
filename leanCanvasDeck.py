import openai
import sys
import os
import json
from pptx import Presentation

# Get the GPT-3.5-turbo API key from the OPENAI_API_KEY environment variable
api_key = os.environ.get('OPENAI_API_KEY')

if not api_key:
    print("Please set the OPENAI_API_KEY environment variable.")
    sys.exit(1)

# Set up the OpenAI API client
openai.api_key = api_key

# Function to generate text using GPT-3.5-turbo
def generate_text(prompt):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=2048,
        temperature=0.7,
        n=1,
        stop=None
    )
    return response.choices[0].text.strip()

# Get the idea description from the command line argument
if len(sys.argv) < 2:
    print("Please provide the description of the idea as a command line argument.")
    sys.exit(1)

idea_description = sys.argv[1]

# Generate the contents of the Lean Canvas for the idea in JSON format
lean_canvas_prompt = f"Create a lean canvas for this business idea: {idea_description}\n\nYour response must be in JSON format like: {{ \"Problem\": [\"foo\", \"bar\"] }}\n\nTry to add 3 bullet points in each category.\n\nYour response:\n\nAI-RESPONSE:"
lean_canvas_text = generate_text(lean_canvas_prompt)

# Parse the JSON response to get the Lean Canvas data
lean_canvas = json.loads(lean_canvas_text)

# Generate the fancy name and tagline for the idea in JSON format
name_and_tagline_prompt = f"Generate a fancy name and tag line for this idea: {idea_description}\n\nYour response must be in JSON format like: {{ \"name\": \"bla\", \"tagline\": \"foo\" }}\n\nYour response:\n\nAI-RESPONSE:"
name_and_tagline_text = generate_text(name_and_tagline_prompt)

# Parse the JSON response to get the fancy name and tagline
name_and_tagline = json.loads(name_and_tagline_text)
fancy_name = name_and_tagline["name"]
tagline = name_and_tagline["tagline"]

# Create a new PowerPoint presentation
presentation = Presentation()

# Add the first slide with the fancy name and tagline
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title_placeholder = slide.placeholders[0]
subtitle_placeholder = slide.placeholders[1]
title_placeholder.text = fancy_name
subtitle_placeholder.text = tagline

# Loop through the Lean Canvas and create a slide for each section
for section, content in lean_canvas.items():
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.placeholders[0]
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = section
    content_placeholder.text = '\n'.join(content)

# Save the presentation to a file
presentation.save(f'{fancy_name}.pptx')

print(f'PowerPoint presentation "{fancy_name}.pptx" has been created successfully.')
